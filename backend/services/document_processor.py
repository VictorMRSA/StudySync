"""
Serviço de processamento de documentos.
Suporta PDF, DOCX, PPTX e TXT.
"""

import io
import logging

from fastapi import UploadFile
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Processador de documentos para extração de texto."""

    # Mapeamento de extensões / content-types para parsers
    _CONTENT_TYPE_MAP: dict[str, str] = {
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "text/plain": "txt",
    }

    _EXTENSION_MAP: dict[str, str] = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".pptx": "pptx",
        ".txt": "txt",
    }

    # ── Parsers individuais ──────────────────────────────────────────────

    async def parse_pdf(self, file: UploadFile) -> str:
        """Extrai texto de todas as páginas de um arquivo PDF."""
        content = await file.read()
        reader = PdfReader(io.BytesIO(content))
        pages_text: list[str] = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages_text.append(text.strip())
        return "\n\n".join(pages_text)

    async def parse_docx(self, file: UploadFile) -> str:
        """Extrai texto de todos os parágrafos de um arquivo DOCX."""
        content = await file.read()
        doc = DocxDocument(io.BytesIO(content))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    async def parse_pptx(self, file: UploadFile) -> str:
        """Extrai texto de todos os slides de um arquivo PPTX."""
        content = await file.read()
        prs = Presentation(io.BytesIO(content))
        slides_text: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slide_content = " ".join(texts)
                slides_text.append(f"\n--- Slide {idx} ---\n{slide_content}")
        return "\n\n".join(slides_text)

    async def parse_txt(self, file: UploadFile) -> str:
        """Lê arquivo de texto puro em UTF-8."""
        content = await file.read()
        return content.decode("utf-8")

    # ── Roteador principal ───────────────────────────────────────────────

    def _detect_format(self, file: UploadFile) -> str | None:
        """Detecta o formato do arquivo pelo content_type ou extensão."""
        # Tentar pelo content type primeiro
        if file.content_type and file.content_type in self._CONTENT_TYPE_MAP:
            return self._CONTENT_TYPE_MAP[file.content_type]

        # Fallback: extensão do nome do arquivo
        if file.filename:
            for ext, fmt in self._EXTENSION_MAP.items():
                if file.filename.lower().endswith(ext):
                    return fmt

        return None

    async def process(self, file: UploadFile) -> tuple[str, str]:
        """
        Processa um arquivo enviado e retorna (texto_extraído, nome_do_arquivo).

        Raises:
            ValueError: se o formato não for suportado.
        """
        file_format = self._detect_format(file)
        filename = file.filename or "documento"

        if file_format is None:
            supported = ", ".join(sorted(self._EXTENSION_MAP.keys()))
            raise ValueError(
                f"Formato não suportado: {file.content_type}. "
                f"Use arquivos {supported}."
            )

        logger.info("Processando arquivo: %s (formato: %s)", filename, file_format)

        parser_map = {
            "pdf": self.parse_pdf,
            "docx": self.parse_docx,
            "pptx": self.parse_pptx,
            "txt": self.parse_txt,
        }

        text = await parser_map[file_format](file)

        if not text or not text.strip():
            raise ValueError(
                "Não foi possível extrair texto do documento. "
                "Verifique se o arquivo contém texto legível."
            )

        logger.info(
            "Texto extraído com sucesso: %d caracteres de '%s'",
            len(text),
            filename,
        )
        return text, filename


# Instância global reutilizável
document_processor = DocumentProcessor()
